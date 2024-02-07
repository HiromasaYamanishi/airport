from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import numpy as np
from pptx.dml.color import RGBColor

color_dict = {
    'black': RGBColor(0, 0, 0),
    'red': RGBColor(255, 0, 0),
    'darkred': RGBColor(139, 0, 0),
    'darkblue': RGBColor(0, 0, 139),
    'green': RGBColor(0, 255, 0),
    'blue': RGBColor(0, 0, 255),
    # 他の色も必要に応じて追加
}

# 色名からRGBColorオブジェクトを取得する関数
def get_rgb_color(color_name):
    return color_dict.get(color_name.lower(), RGBColor(0, 0, 0)) 

class SlideWriter:
    def __init__(self, ):
        self.prs = Presentation()
        
    def write_tables(self, arrays:np.array, row_names, col_names, colors=np.array([]), font_size=15):
        '''
        arrays: 載せる表のリスト
        row_names: 各表における列名のリスト
        col_names: 各表における行名のリスト
        '''
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # 空のレイアウトを選択

        # 最初の表の開始位置
        y_offset = Inches(1.5)

        for array, rows, cols in zip(arrays, row_names, col_names):
            # 表の行数と列数を計算
            num_rows = len(array) + 1  # データ行数に列名行を追加
            num_cols = len(cols) + 1  # データ列数に行名列を追加
            print(num_rows, num_cols)
            # 表をスライドに追加
            table = slide.shapes.add_table(num_rows, num_cols, Inches(1), y_offset, Inches(10), Inches(1 * num_rows)).table

            # 列名を設定
            for i, col_name in enumerate(cols):
                table.cell(0, i+1).text = col_name

            # 行名とデータを挿入
            for row_idx, (row_name, row_data) in enumerate(zip(rows, array), start=1):
                table.cell(row_idx, 0).text = row_name  # 行名
                for col_idx, item in enumerate(row_data, start=1):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(item)
                    table.cell(row_idx, col_idx).text = str(item)
                    text_frame = cell.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)  # フォントサイズを12ポイントに設定
                            run.font.name = "Calibri"
                            if colors.shape==array.shape:
                                run.font.color.rgb = get_rgb_color(colors[row_idx-1, col_idx-1])
            y_offset += Inches(0.5 * num_rows + 0.5)

    def save(self, filename):
        self.prs.save(filename)

if __name__=='__main__':
    slide_writer = SlideWriter()

    # テストデータの作成
    arrays = [
        [[f"データ1-{r+1}-{c+1}" for c in range(5)] for r in range(3)],
        [[f"データ2-{r+1}-{c+1}" for c in range(4)] for r in range(4)]
    ]
    row_names = [
        ["行1", "行2", "行3"],
        ["行A", "行B", "行C", "行D"]
    ]
    col_names = [
        ["列1", "列2", "列3", "列4", "列5"],
        ["列A", "列B", "列C", "列D"]
    ]

    slide_writer.write_tables(arrays, row_names, col_names)
    slide_writer.save('updated_presentation.pptx')